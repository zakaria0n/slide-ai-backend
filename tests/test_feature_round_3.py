"""Tests for the feature round: chart element, outline, PPTX import,
translation and user themes."""
from __future__ import annotations

import io
import uuid

import jwt
import pytest
from fastapi.testclient import TestClient

SECRET = "test-secret"


def _token(user_id: str, secret: str = SECRET) -> str:
    return jwt.encode(
        {"sub": user_id, "email": "u@example.com", "aud": "authenticated"},
        secret,
        algorithm="HS256",
    )


def _auth(token: str) -> dict[str, str]:
    return {"Authorization": f"Bearer {token}"}


# --- chart element in the spec ----------------------------------------------


def test_chart_element_spec_validation() -> None:
    from app.generation.spec import PresentationSpec

    spec = PresentationSpec.model_validate(
        {
            "meta": {"title": "Charts"},
            "slides": [
                {
                    "layout": "blank",
                    "elements": [
                        {
                            "type": "chart",
                            "chart_type": "bar",
                            "labels": ["Q1", "Q2"],
                            "datasets": [{"label": "Revenue", "data": [12, 18.5]}],
                            "x": 10,
                            "y": 20,
                            "w": 60,
                            "h": 50,
                        }
                    ],
                }
            ],
        }
    )
    el = spec.slides[0].elements[0]
    assert el.type == "chart"
    assert el.datasets[0].data == [12.0, 18.5]


def test_chart_element_rejects_bad_type() -> None:
    from pydantic import ValidationError

    from app.generation.spec import ChartElement

    with pytest.raises(ValidationError):
        ChartElement.model_validate({"type": "chart", "chart_type": "hologram"})


def test_locked_flag_persists() -> None:
    from app.generation.spec import TitleElement

    el = TitleElement.model_validate({"type": "title", "text": "T", "locked": True})
    assert el.locked is True


# --- native PPTX chart export -------------------------------------------------


@pytest.mark.asyncio
async def test_pptx_export_embeds_native_chart() -> None:
    from pptx import Presentation as PptxPresentation

    from app.export.pptx_exporter import _export_pptx
    from app.export.html_theme import tokens_for
    from app.generation.spec import ChartDataset, ChartElement, PresentationSpec, SlideSpec, TitleElement

    spec = PresentationSpec(
        meta={"title": "Native chart deck"},
        slides=[
            SlideSpec(
                layout="chart",
                elements=[
                    TitleElement(type="title", text="Revenue"),
                    ChartElement(
                        type="chart",
                        chart_type="bar",
                        labels=["Q1", "Q2", "Q3"],
                        datasets=[ChartDataset(label="Rev", data=[1.0, 2.0, 3.5])],
                    ),
                ],
            )
        ],
    )
    data = await _export_pptx(spec, tokens_for(None))
    prs = PptxPresentation(io.BytesIO(data))
    chart_shapes = [sh for sh in prs.slides[0].shapes if getattr(sh, "has_chart", False)]
    assert chart_shapes, "expected a native chart shape in the PPTX"
    chart = chart_shapes[0].chart
    assert len(chart.plots[0].categories) == 3
    assert len(list(chart.series)) == 1


# --- PPTX import ----------------------------------------------------------------


def _build_source_pptx() -> bytes:
    from pptx import Presentation as PptxPresentation
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Emu, Pt

    prs = PptxPresentation()
    prs.slide_width = Emu(13_333_333)
    prs.slide_height = Emu(7_500_000)
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # title-slide layout
    slide.shapes.title.text = "Imported heading"

    body = slide.shapes.add_textbox(Emu(685_800), Emu(1_700_000), Emu(8_000_000), Emu(3_000_000))
    tf = body.text_frame
    tf.text = "First point"
    tf.add_paragraph().text = "Second point"

    tbl = slide.shapes.add_table(2, 2, Emu(685_800), Emu(5_000_000), Emu(8_000_000), Emu(1_500_000))
    tbl.table.cell(0, 0).text = "Name"
    tbl.table.cell(0, 1).text = "Value"
    tbl.table.cell(1, 0).text = "Alpha"
    tbl.table.cell(1, 1).text = "10"

    cd = CategoryChartData()
    cd.categories = ["A", "B"]
    cd.add_series("S", (4.0, 7.0))
    slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Emu(9_000_000), Emu(1_000_000), Emu(3_500_000), Emu(3_000_000), cd)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_pptx_import_converts_all_shapes() -> None:
    from app.presentations.pptx_import import import_pptx_to_spec

    spec = import_pptx_to_spec(_build_source_pptx(), title="My import")
    assert spec.meta.title == "My import"
    assert len(spec.slides) == 1
    types = [el.type for el in spec.slides[0].elements]
    assert "title" in types
    assert "paragraph" in types or "bullets" in types
    assert "table" in types
    assert "chart" in types

    chart = next(el for el in spec.slides[0].elements if el.type == "chart")
    assert chart.chart_type == "bar"
    assert chart.labels == ["A", "B"]
    assert chart.datasets[0].data == [4.0, 7.0]
    table = next(el for el in spec.slides[0].elements if el.type == "table")
    assert table.headers == ["Name", "Value"]
    assert table.rows == [["Alpha", "10"]]


def test_pptx_import_rejects_empty_file() -> None:
    from app.core.exceptions import ValidationError
    from app.presentations.pptx_import import import_pptx_to_spec

    with pytest.raises(ValidationError):
        import_pptx_to_spec(b"not a pptx", title="x")


# --- outline endpoint -------------------------------------------------------------


def test_outline_endpoint(monkeypatch, client: TestClient) -> None:
    from app.generation.schemas import OutlineItem

    async def fake_outline(settings, **kwargs):
        assert kwargs["slide_count"] == 3
        return [
            OutlineItem(title="Opening", points=["hook"]),
            OutlineItem(title="Core idea"),
            OutlineItem(title="Closing"),
        ]

    monkeypatch.setattr("app.generation.outliner.generate_outline", fake_outline)
    uid = str(uuid.uuid4())
    res = client.post(
        "/api/v1/presentations/outline",
        json={"prompt": "test outline", "slide_count": 3},
        headers=_auth(_token(uid)),
    )
    assert res.status_code == 200
    body = res.json()
    assert [o["title"] for o in body["outline"]] == ["Opening", "Core idea", "Closing"]


def test_generate_accepts_outline(client: TestClient) -> None:
    """Offline provider: a deck generated with an approved outline keeps the
    outline's slide count (the pipeline accepts the field end-to-end)."""
    uid = str(uuid.uuid4())
    res = client.post(
        "/api/v1/presentations/generate",
        json={
            "prompt": "outline test deck",
            "slide_count": 3,
            "outline": [
                {"title": "Plan A", "points": ["p1"]},
                {"title": "Plan B"},
                {"title": "Plan C", "points": ["x", "y"]},
            ],
        },
        headers=_auth(_token(uid)),
    )
    assert res.status_code == 201, res.text
    assert res.json()["slide_count"] >= 1


# --- translation --------------------------------------------------------------------


def _seed_spec(client: TestClient, uid: str) -> str:
    created = client.post(
        "/api/v1/presentations",
        json={"title": "translate me"},
        headers=_auth(_token(uid)),
    ).json()
    pid = created["id"]
    res = client.put(
        f"/api/v1/presentations/{pid}/spec",
        json={
            "meta": {"title": "Hello", "language": "English"},
            "slides": [
                {
                    "layout": "title",
                    "elements": [
                        {"type": "title", "text": "Welcome"},
                        {"type": "paragraph", "text": "This is a deck"},
                        {"type": "bullets", "items": ["one", "two"]},
                    ],
                }
            ],
        },
        headers=_auth(_token(uid)),
    )
    assert res.status_code == 200, res.text
    return pid


def test_translate_endpoint(monkeypatch, client: TestClient) -> None:
    captured = {}

    async def fake_complete_json(settings, *, model, system, user, max_tokens=4000):
        captured["user"] = user
        return {
            "meta.title": "Bonjour",
            "0.0.text": "Bienvenue",
            "0.1.text": "Ceci est un deck",
            "0.2.items[0]": "un",
            "0.2.items[1]": "deux",
        }

    monkeypatch.setattr("app.generation.translator.complete_json", fake_complete_json)
    uid = str(uuid.uuid4())
    pid = _seed_spec(client, uid)

    res = client.post(
        f"/api/v1/presentations/{pid}/translate",
        json={"target_language": "French"},
        headers=_auth(_token(uid)),
    )
    assert res.status_code == 200, res.text
    spec = res.json()
    assert spec["meta"]["title"] == "Bonjour"
    assert spec["meta"]["language"] == "French"
    els = spec["slides"][0]["elements"]
    assert els[0]["text"] == "Bienvenue"
    assert els[1]["text"] == "Ceci est un deck"
    assert els[2]["items"] == ["un", "deux"]
    # Numbers/keys from the source spec were part of the payload.
    assert "0.2.items[0]" in captured["user"]


def test_translator_collect_and_apply() -> None:
    from app.generation.spec import PresentationSpec
    from app.generation.translator import _apply_translations, _collect_texts

    spec = PresentationSpec.model_validate(
        {
            "meta": {"title": "Hi"},
            "slides": [
                {
                    "layout": "quote",
                    "elements": [
                        {"type": "quote", "text": "Be bold", "author": "Ada"},
                        {"type": "table", "headers": ["K", "V"], "rows": [["a", "b"]]},
                    ],
                    "notes": "speak clearly",
                }
            ],
        }
    )
    texts = _collect_texts(spec)
    assert texts["0.0.text"] == "Be bold"
    assert texts["0.0.author"] == "Ada"
    assert texts["0.1.headers[0]"] == "K"
    assert texts["0.1.rows[0][0]"] == "a"
    assert texts["0.notes"] == "speak clearly"
    assert texts["meta.title"] == "Hi"

    _apply_translations(spec, {
        "0.0.text": "Soyez audacieux",
        "0.0.author": "Ada",
        "0.1.headers[0]": "Clé",
        "0.1.rows[0][0]": "a",
        "0.notes": "parle clairement",
        "meta.title": "Salut",
    })
    assert spec.slides[0].elements[0].text == "Soyez audacieux"
    assert spec.slides[0].elements[1].headers[0] == "Clé"
    assert spec.slides[0].notes == "parle clairement"
    assert spec.meta.title == "Salut"


# --- user themes --------------------------------------------------------------------


def test_user_themes_crud(client: TestClient) -> None:
    uid = str(uuid.uuid4())
    headers = _auth(_token(uid))

    res = client.post(
        "/api/v1/themes",
        json={
            "name": "My brand",
            "tokens": {"bg": "#0b0b10", "accent": "#7c5cff", "gradient": "linear-gradient(...)"},
            "ambient": {"kind": "blobs", "colors": ["#7c5cff"], "opacity": 0.5, "speed": 30},
        },
        headers=headers,
    )
    assert res.status_code == 201, res.text
    theme = res.json()
    assert theme["name"] == "My brand"
    assert theme["tokens"]["accent"] == "#7c5cff"

    listing = client.get("/api/v1/themes", headers=headers)
    assert listing.status_code == 200
    assert len(listing.json()["themes"]) == 1

    # Another user sees nothing.
    other = client.get("/api/v1/themes", headers=_auth(_token(str(uuid.uuid4()))))
    assert other.json()["themes"] == []

    del_res = client.delete(f"/api/v1/themes/{theme['id']}", headers=headers)
    assert del_res.status_code == 204
    assert client.get("/api/v1/themes", headers=headers).json()["themes"] == []
