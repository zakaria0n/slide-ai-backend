"""PPTX export strategy.

Generates a native PowerPoint file with the deck's *content only* — no
animations (matching the Phase 12 spec). One slide per spec slide. Uses
python-pptx; the brand is surfaced as "Slide AI" only via the document
title/property, never the underlying AI provider.
"""
from __future__ import annotations

from io import BytesIO
from typing import Any

from pptx import Presentation as PptxPresentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from app.export.html_theme import ThemeTokens, tokens_for
from app.export.strategy import ExportFormat, ExportStrategy, ExportedFile
from app.generation.spec import PresentationSpec


def _hex(color: str) -> RGBColor:
    c = color.lstrip("#")
    if len(c) == 3:
        c = "".join(ch * 2 for ch in c)
    try:
        return RGBColor.from_string(c[:6])
    except Exception:
        return RGBColor(0x10, 0x10, 0x16)


def _group(slide: dict) -> dict[str, list[dict]]:
    by_type: dict[str, list[dict]] = {}
    for el in slide.get("elements", []):
        by_type.setdefault(el.get("type", ""), []).append(el)
    return by_type


def _add_text(tf, text: str, size: int, color: RGBColor, bold: bool = False, italic: bool = False):
    p = tf.paragraphs[0] if tf.paragraphs[0].text == "" else tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    return p


def _export_pptx(spec: PresentationSpec, t: ThemeTokens) -> bytes:
    prs = PptxPresentation()
    prs.slide_width = Emu(13_333_333)  # 16:9
    prs.slide_height = Emu(7_500_000)
    blank = prs.slide_layouts[6]
    bg = _hex(t.bg)
    accent = _hex(t.accent)
    accent2 = _hex(t.accent2)
    text_c = _hex(t.text)
    muted_c = _hex(t.text_muted)

    for slide in spec.slides:
        s = slide.model_dump() if hasattr(slide, "model_dump") else slide
        ppt_slide = prs.slides.add_slide(blank)
        # background fill
        ppt_slide.background.fill.solid()
        ppt_slide.background.fill.fore_color.rgb = bg
        g = _group(s)

        # Title
        title_box = ppt_slide.shapes.add_textbox(Emu(685_800), Emu(457_200), Emu(11_961_600), Emu(1_200_000))
        tf = title_box.text_frame
        tf.word_wrap = True
        for i, el in enumerate(g.get("title", [])):
            _add_text(tf, str(el.get("text", "")), 40, text_c, bold=True)
        # Subtitle
        for el in g.get("subtitle", []):
            _add_text(tf, str(el.get("text", "")), 24, muted_c)

        # Bullets
        if g.get("bullets"):
            b = g["bullets"][0]
            box = ppt_slide.shapes.add_textbox(Emu(685_800), Emu(1_900_000), Emu(11_961_600), Emu(4_500_000))
            tf = box.text_frame
            tf.word_wrap = True
            first = True
            for item in (b.get("items") or []):
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                p.text = f"• {item}"
                p.font.size = Pt(22)
                p.font.color.rgb = text_c

        # Cards / Statistics as a simple grid of text boxes
        for key in ("cards", "statistics"):
            if g.get(key):
                items = (g[key][0].get("items") or []) if g.get(key) else []
                cols = min(4, max(1, len(items))) if items else 1
                cell_w = int(11_961_600 / cols)
                for idx, it in enumerate(items):
                    cx = Emu(685_800 + cell_w * (idx % cols))
                    cy = Emu(2_200_000 + (idx // cols) * 1_600_000)
                    box = ppt_slide.shapes.add_textbox(cx, cy, Emu(cell_w - 200_000), Emu(1_400_000))
                    tf = box.text_frame
                    tf.word_wrap = True
                    _add_text(tf, str(it.get("title", it.get("value", ""))), 20, accent, bold=True)
                    sub = it.get("body", it.get("label", ""))
                    if sub:
                        _add_text(tf, str(sub), 14, muted_c)

        # Comparison
        if g.get("comparison"):
            cmp = g["comparison"][0]
            for ci, col in enumerate((cmp.get("left", {}), cmp.get("right", {}))):
                cx = Emu(685_800 + 6_000_000 * ci)
                box = ppt_slide.shapes.add_textbox(cx, Emu(2_000_000), Emu(5_800_000), Emu(4_500_000))
                tf = box.text_frame
                tf.word_wrap = True
                _add_text(tf, str(col.get("title", "")), 22, accent2 if ci else accent, bold=True)
                for pnt in (col.get("points") or []):
                    _add_text(tf, f"• {pnt}", 16, text_c)

        # Quote
        for el in g.get("quote", []):
            box = ppt_slide.shapes.add_textbox(Emu(1_500_000), Emu(2_500_000), Emu(10_300_000), Emu(3_000_000))
            tf = box.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            _add_text(tf, f"“{el.get('text','')}”", 28, text_c, italic=True)
            if el.get("author"):
                _add_text(tf, f"— {el.get('author')}", 18, muted_c)

    # Document property (brand only)
    prs.core_properties.title = (getattr(spec.meta, "title", "Slide AI Presentation") if spec.meta else "Slide AI Presentation")
    prs.core_properties.author = "Slide AI"

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


class PptxExportStrategy(ExportStrategy):
    format = ExportFormat.PPTX

    def export(self, spec: PresentationSpec, theme_hint: str | None = None) -> ExportedFile:
        t = tokens_for(theme_hint or (getattr(spec.meta, "theme", None) if spec.meta else None))
        data = _export_pptx(spec, t)
        title = getattr(spec.meta, "title", "presentation") if spec.meta else "presentation"
        safe = "".join(c if c.isalnum() else "-" for c in str(title)).strip("-") or "presentation"
        return ExportedFile(data, "application/vnd.openxmlformats-officedocument.presentationml.presentation", f"{safe}.pptx")
