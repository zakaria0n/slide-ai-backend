"""Import an existing .pptx file into a Slide AI presentation spec.

Deterministic conversion (no AI): each PowerPoint slide becomes a spec
slide whose elements keep their free positions, so the imported deck is
editable on the canvas right away. Text placeholders become title /
paragraph / bullets elements, tables become table elements and native
charts become REAL chart elements (chart type + labels + series data).
Pictures are skipped (binary assets are not copied); speaker notes are
imported.
"""
from __future__ import annotations

from io import BytesIO

from pptx import Presentation as PptxPresentation

from app.generation.spec import (
    BulletsElement,
    ChartDataset,
    ChartElement,
    ParagraphElement,
    PresentationSpec,
    SlideSpec,
    TableElement,
    TitleElement,
)

_EMU_PER_INCH = 914_400


def _pct(value: int | None, total: int | None, default: float) -> float:
    """EMU length → percent of the slide dimension, clamped 0-100."""
    if not value or not total:
        return default
    return max(0.0, min(100.0, round(value / total * 100, 2)))


def _paragraph_text(paragraph) -> str:
    return "".join(run.text for run in paragraph.runs).strip()


def _title_text(slide) -> str:
    try:
        if slide.shapes.title is not None:
            return (slide.shapes.title.text or "").strip()
    except Exception:
        pass
    return ""


def _convert_chart(shape) -> ChartElement | None:
    """Read a native PowerPoint chart into a chart element."""
    try:
        chart = shape.chart
        plots = chart.plots
        if not plots:
            return None
        plot = plots[0]
        labels = [str(c) for c in plot.categories]
        datasets = []
        for series in plot.series:
            datasets.append(
                ChartDataset(
                    label=str(series.name or ""),
                    data=[float(v) if v is not None else 0.0 for v in series.values],
                )
            )
        if not datasets:
            return None
        kind_map = {
            "BAR": "bar",
            "COLUMN": "bar",
            "LINE": "line",
            "PIE": "pie",
            "DOUGHNUT": "doughnut",
            "RADAR": "radar",
        }
        # XL_CHART_TYPE enum member name, e.g. "COLUMN_CLUSTERED" / "PIE".
        raw_kind = getattr(chart.chart_type, "name", str(chart.chart_type))
        kind = kind_map.get(str(raw_kind).split("_")[0].upper(), "bar")
        return ChartElement(
            type="chart",
            chart_type=kind,
            labels=labels[:24],
            datasets=datasets[:6],
        )
    except Exception:
        return None


def _convert_table(shape) -> TableElement | None:
    try:
        table = shape.table
        rows = [[str(cell.text or "").strip() for cell in row.cells] for row in table.rows]
        if not rows:
            return None
        return TableElement(
            type="table",
            headers=rows[0],
            rows=[r for r in rows[1:] if any(r)],
        )
    except Exception:
        return None


def _convert_text_shape(shape, is_body: bool):
    """Text frame → paragraph (single block) or bullets (multi-line body)."""
    try:
        tf = shape.text_frame
    except Exception:
        return None
    lines = [_paragraph_text(p) for p in tf.paragraphs]
    lines = [ln for ln in lines if ln]
    if not lines:
        return None
    if is_body and len(lines) > 1:
        return BulletsElement(type="bullets", items=lines[:20])
    return ParagraphElement(type="paragraph", text="\n".join(lines[:30]))


def _placeholder_kind(shape) -> bool:
    """True when the shape is a body/object content placeholder."""
    try:
        ph = shape.placeholder_format
        return ph is not None and ph.type is not None and "PICTURE" not in str(ph.type)
    except Exception:
        return False


def import_pptx_to_spec(data: bytes, *, title: str = "", theme: str | None = None) -> PresentationSpec:
    """Convert a .pptx byte stream into a PresentationSpec."""
    from app.core.exceptions import ValidationError

    try:
        prs = PptxPresentation(BytesIO(data))
    except Exception as exc:
        raise ValidationError("This file is not a valid .pptx presentation") from exc
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    spec_slides: list[SlideSpec] = []
    for slide in prs.slides:
        elements = []
        title_text = _title_text(slide)
        if title_text:
            elements.append(TitleElement(type="title", text=title_text[:300]))
        # Identity of the title shape so its text is not exported twice.
        try:
            title_shape_id = slide.shapes.title.shape_id if slide.shapes.title else None
        except Exception:
            title_shape_id = None

        for shape in slide.shapes:
            if title_shape_id is not None and shape.shape_id == title_shape_id:
                continue

            x = _pct(getattr(shape, "left", None), slide_w, 6.0)
            y = _pct(getattr(shape, "top", None), slide_h, 20.0)
            w = _pct(getattr(shape, "width", None), slide_w, 40.0)
            h = _pct(getattr(shape, "height", None), slide_h, 16.0)

            # Charts (GraphicFrame with a chart part).
            if getattr(shape, "has_chart", False):
                chart_el = _convert_chart(shape)
                if chart_el is not None:
                    chart_el.x, chart_el.y, chart_el.w, chart_el.h = x, y, max(w, 30.0), max(h, 25.0)
                    elements.append(chart_el)
                continue

            # Tables.
            if getattr(shape, "has_table", False):
                tbl = _convert_table(shape)
                if tbl is not None:
                    tbl.x, tbl.y, tbl.w, tbl.h = x, y, max(w, 50.0), max(h, 20.0)
                    elements.append(tbl)
                continue

            # Plain text (skip picture placeholders — binary assets are not copied).
            if shape.has_text_frame:
                el = _convert_text_shape(shape, is_body=_placeholder_kind(shape))
                if el is not None:
                    el.x, el.y, el.w, el.h = x, y, max(w, 10.0), max(h, 6.0)
                    elements.append(el)

        notes = ""
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
                notes = (slide.notes_slide.notes_text_frame.text or "").strip()[:2000]
        except Exception:
            notes = ""

        spec_slides.append(
            SlideSpec(
                layout="blank",
                elements=elements,
                notes=notes or None,
            )
        )

    if not spec_slides:
        raise ValidationError("The .pptx file contains no slides")

    meta_title = (title or "").strip() or "Imported presentation"
    return PresentationSpec(
        meta={"title": meta_title, "theme": theme},
        slides=spec_slides,
    )
